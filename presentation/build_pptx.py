#!/usr/bin/env python3
"""Generate a 32-slide RoommateSync .pptx presentation with embedded Mermaid diagrams."""
import os, textwrap, subprocess, html
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = r"D:\ISD_PROJECT\ISD-Project-RoommateSync"
OUT_DIR = os.path.join(ROOT, "presentation")
ASSET_DIR = os.path.join(OUT_DIR, "diagrams")
os.makedirs(ASSET_DIR, exist_ok=True)
PPTX_PATH = os.path.join(OUT_DIR, "RoommateSync-Presentation.pptx")
MMDC = r"C:\Users\USER\AppData\Roaming\npm\mmdc.cmd"

# ---- palette ----
BG      = RGBColor(0x0C, 0x0F, 0x1A)
PANEL   = RGBColor(0x14, 0x1A, 0x2E)
PANEL2  = RGBColor(0x1B, 0x22, 0x36)
ACCENT  = RGBColor(0x5B, 0x8C, 0xFF)
ACCENT2 = RGBColor(0x7C, 0x5C, 0xFF)
TEXT    = RGBColor(0xE8, 0xED, 0xF7)
MUTED   = RGBColor(0x9A, 0xA6, 0xC2)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OK      = RGBColor(0x3D, 0xDC, 0x97)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

mmd_count = 0
def add_diagram(slide, code, left, top, width, height):
    """Render a mermaid block to PNG and place it on the slide."""
    global mmd_count
    mmd_count += 1
    base = os.path.join(ASSET_DIR, f"d{mmd_count}")
    mmd_path = base + ".mmd"
    png_path = base + ".png"
    with open(mmd_path, "w", encoding="utf-8") as f:
        f.write(code)
    try:
        r = subprocess.run(
            [MMDC, "-i", mmd_path, "-o", png_path, "-t", "dark",
             "-b", "transparent", "--width", "1200", "--scale", "2"],
            check=True, capture_output=True, text=True
        )
    except subprocess.CalledProcessError as e:
        print(f"mmdc FAILED on diagram #{mmd_count}")
        print(e.stderr[:400])
        return
    slide.shapes.add_picture(png_path, left, top, width=width)

def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=PANEL, line=RGBColor(0x2A,0x33,0x4D)):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(para, tuple):
            para = [para]
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Calibri"
    if fill is not None:
        box(slide, l, t, w, h, fill=fill)
    return tb

def kicker(slide, txt, l=Inches(0.6), t=Inches(0.45)):
    text(slide, l, t, Inches(8), Inches(0.35),
         [[(txt.upper(), 13, ACCENT, True)]])

def heading(slide, txt, size=32, l=Inches(0.6), t=Inches(0.8)):
    text(slide, l, t, Inches(12.1), Inches(0.9),
         [[(txt, size, WHITE, True)]])

def bullets(slide, l, t, w, h, items, size=18, gap=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            lead, rest = it
            r = p.add_run(); r.text = "• " + lead
            r.font.size = Pt(size); r.font.color.rgb = ACCENT; r.font.bold = True; r.font.name="Calibri"
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = MUTED; r2.font.name="Calibri"
        else:
            r = p.add_run(); r.text = "• " + it
            r.font.size = Pt(size); r.font.color.rgb = MUTED; r.font.name="Calibri"
    return tb

def table(slide, l, t, w, rows, col_w=None, header_fill=ACCENT, fs=15):
    n = len(rows[0])
    gt = slide.shapes.add_table(len(rows), n, l, t, w, Inches(0.4*len(rows))).table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            gt.columns[i].width = Emu(int(w * cw / total))
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fs); r.font.name = "Calibri"
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                r.font.bold = True; r.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL2 if ri % 2 else PANEL
                r.font.color.rgb = MUTED if ci != 0 else TEXT
    return gt

# ============ SLIDES ============
# 1 TITLE
s = prs.slides.add_slide(BLANK); bg(s)
text(s, Inches(2), Inches(2.2), Inches(9.3), Inches(1.0),
     [[("ROOMMATESYNC", 60, WHITE, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(2), Inches(3.4), Inches(9.3), Inches(0.8),
     [[("Find a room, split costs, manage house life — all in one place.", 24, ACCENT, False)]],
     align=PP_ALIGN.CENTER)
pills = ["PHP 8.1","MySQL 8.0","HTML5 / CSS3","Vanilla JS","XAMPP","MIT"]
text(s, Inches(2), Inches(4.6), Inches(9.3), Inches(0.6),
     [[("  ".join("● "+p for p in pills), 14, MUTED, False)]], align=PP_ALIGN.CENTER)
text(s, Inches(2), Inches(5.6), Inches(9.3), Inches(0.5),
     [[("Team RoommateSync · Dadhichi · Shawki · Plabon  —  ISD Course Project, KUET", 16, MUTED, False)]],
     align=PP_ALIGN.CENTER)

# 2 AGENDA
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "What we'll cover"); heading(s, "Presentation Agenda")
agenda = [
    ("1. Problem & Solution", " — why RoommateSync exists and what it delivers."),
    ("2. Modules", " — the seven integrated feature modules."),
    ("3. Architecture", " — layers, data flow, deployment."),
    ("4. Data Model", " — ERD, schema, relationships."),
    ("5. UML Diagrams", " — use case, class, sequence, state, activity & more."),
    ("6. Process & Results", " — sprints, testing, team & demo."),
]
cards = [(agenda[i], agenda[i+1]) for i in range(0,6,2)]
for idx,(a,b) in enumerate(cards):
    col = idx % 2; rowi = idx // 2
    l = Inches(0.6 + col*6.35); t = Inches(1.9 + rowi*1.7)
    box(s, l, t, Inches(6.0), Inches(1.5))
    bullets(s, l+Inches(0.2), t+Inches(0.15), Inches(5.6), Inches(1.2), [a,b], size=15, gap=2)

# 3 PROBLEM
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "The Context"); heading(s, "The Problem")
text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.9),
     [[("Students and young professionals relocating to cities like Dhaka struggle to find trustworthy "
        "roommates, verify listings, and manage shared expenses.", 20, TEXT, False)]])
probs = [
    ("Fragmented discovery", " — Facebook groups, word-of-mouth and spreadsheets, no single source of truth."),
    ("Trust gap", " — no way to verify landlords or rate past roommates on compatibility."),
    ("Shared-cost chaos", " — splitting bills is manual, unfair and easy to lose track of."),
]
for i,(lead,rest) in enumerate(probs):
    l = Inches(0.6 + i*4.15); t = Inches(3.0)
    box(s, l, t, Inches(3.9), Inches(2.6))
    bullets(s, l+Inches(0.2), t+Inches(0.2), Inches(3.5), Inches(2.2), [(lead,rest)], size=15, gap=4)
text(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.6),
     [[("Built as a university ISD project by a team of 3 students for the Bangladeshi rental market.", 16, MUTED, False)]])

# 4 SOLUTION
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Our Approach"); heading(s, "The Solution — One Unified Platform")
text(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5),
     [[("A full-stack PHP web app that unifies discovery, trust and shared-living management.", 19, TEXT, False)]])
rows = [
    ["#","Module","What it does"],
    ["1","Rental Marketplace","Browse, filter & search listings by price, room type & location — real-time."],
    ["2","Bill Split Calculator","Split bills proportionally by income with optional DB logging."],
    ["3","Viewing Booking","30-min slots with automatic conflict prevention."],
    ["4","Listing Upload","Landlords create listings with image upload & validation."],
    ["5","Peer Review","Rate roommates; aggregated cleanliness & communication scores."],
    ["6","Chat","Real-time polling messaging between accepted connections."],
    ["7","Connect","Double opt-in matching that unlocks chat when both accept."],
]
table(s, Inches(0.6), Inches(2.2), Inches(12.1), rows, col_w=[0.5,2.4,8.5], fs=14)

# 5 TECH STACK
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Under the Hood"); heading(s, "Technology Stack")
stack = [
    ("Frontend","HTML5, CSS3, Vanilla JS — glassmorphism dark theme, responsive grid."),
    ("Backend","PHP 8.1 (pure, no framework) — API endpoints & session management."),
    ("Database","MySQL 8.0 via XAMPP — relational storage, indexing, transactions."),
    ("Data Access","PDO singleton — prepared statements, transaction & lock support."),
    ("Dev & VCS","PHP built-in server (php -S), Git + GitHub, Jira, Zephyr."),
    ("Testing","Zephyr for Jira — 49 test cases across 7 cycles, 100% pass."),
]
for i,(lead,rest) in enumerate(stack):
    col = i % 3; rowi = i // 3
    l = Inches(0.6 + col*4.15); t = Inches(2.0 + rowi*2.2)
    box(s, l, t, Inches(3.9), Inches(1.9))
    bullets(s, l+Inches(0.2), t+Inches(0.2), Inches(3.5), Inches(1.5), [(lead,rest)], size=14, gap=3)

# 6 SYSTEM ARCHITECTURE
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Architecture"); heading(s, "System Architecture")
add_diagram(s, """graph TB
  Browser[Web Browser]
  subgraph PL["Presentation"]
    Dashboard[index.php Dashboard]
    Auth[auth pages]
    Modules[modules pages]
    CSS[style.css]
    JS[assets js]
  end
  subgraph AL["Application"]
    Layout[layout.php]
    Helpers[helpers.php]
    AuthCore[auth.php RBAC]
    Bootstrap[bootstrap.php]
  end
  subgraph DL["Data"]
    Database[database.php PDO]
    MySQL[(MySQL 8.0)]
  end
  Uploads[Uploads/ Images]
  Browser -->|HTTP| Dashboard
  Browser -->|HTTP| Auth
  Browser -->|HTTP| Modules
  Dashboard --> Layout
  Modules --> Layout
  Layout --> AuthCore
  AuthCore --> Database
  Modules --> Database
  Database --> MySQL
  Modules -->|File Upload| Uploads""",
  Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.4))

# 7 DFD L0
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Data Flow Diagram"); heading(s, "DFD Level 0 — Context")
add_diagram(s, """graph LR
  User((Tenant / Landlord))
  System[RoommateSync System]
  DB[(MySQL Database)]
  User -->|Browse, Book, Review, Chat| System
  System -->|Listings, Slots, Messages| User
  System <-->|Read / Write| DB""",
  Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.8))

# 8 DFD L1
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Data Flow Diagram"); heading(s, "DFD Level 1 — Major Processes")
add_diagram(s, """graph TB
  U1[Register Login]
  U2[Search Filter]
  U3[Booking]
  U4[Connect]
  U5[Review]
  U6[Bill Split]
  P1[1.0 Auth]
  P2[2.0 Marketplace]
  P3[3.0 Booking]
  P4[4.0 Social]
  P5[5.0 Review]
  P6[6.0 Bill]
  D1[(users)]
  D2[(listings)]
  D3[(appointments)]
  D4[(connection_requests)]
  D5[(messages)]
  D6[(user_reviews)]
  D7[(bill_logs)]
  U1 --> P1
  U2 --> P2
  U3 --> P3
  U4 --> P4
  U5 --> P5
  U6 --> P6
  P1 <--> D1
  P2 <--> D2
  P3 <--> D3
  P4 <--> D4
  P4 <--> D5
  P5 <--> D6
  P6 <--> D7""",
  Inches(0.4), Inches(1.8), Inches(12.5), Inches(5.4))

# 9 DFD L2
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Data Flow Diagram"); heading(s, "DFD Level 2 — Booking Detail")
add_diagram(s, """graph LR
  A[Select listing] --> B[Pick date]
  B --> C[API available_slots]
  C --> D[Query non-cancelled]
  D --> E[Return booked slots]
  E --> F[Pick time slot]
  F --> G[API book_viewing]
  G --> H{Conflict?}
  H -->|No| I[INSERT status=PENDING]
  H -->|Yes| J[409 error]
  I --> K[Success + booking ID]""",
  Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.8))

# 10 ERD
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Data Model"); heading(s, "Entity Relationship Diagram")
add_diagram(s, """erDiagram
  users ||--o{ listings : owns
  users ||--o{ appointments : books
  users ||--o{ connection_requests : sends
  users ||--o{ user_reviews : writes
  users ||--o{ messages : sends
  users ||--o{ bill_logs : creates
  bill_logs ||--o{ bill_log_roommates : contains
  listings ||--o{ appointments : has
  users ||--o{ user_profiles : has
  users ||--o{ user_profile_tags : has""",
  Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.4))

# 11 SCHEMA
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Data Model"); heading(s, "Database Schema")
rows = [
    ["Table","Rows (seed)","Purpose"],
    ["users","5","Accounts with roles"],
    ["user_profiles","5","Lifestyle preferences"],
    ["user_profile_tags","19","Interest tags"],
    ["listings","5","Rental listings"],
    ["appointments","2","Viewing bookings"],
    ["connection_requests","2","Social pairs"],
    ["messages","2","Chat"],
    ["user_reviews","1","Peer reviews"],
    ["bill_logs / bill_log_roommates","0","Saved splits"],
]
table(s, Inches(0.6), Inches(1.9), Inches(12.1), rows, col_w=[4,2,6], fs=14)
text(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
     [[("11 tables · 8 indexes · 12 foreign keys.", 18, ACCENT, True)]])

# 12 USE CASE
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Use Case Diagram")
add_diagram(s, """graph TB
  Tenant((Tenant))
  Landlord((Landlord))
  Admin((Admin))
  UC1[Register Login]
  UC2[Browse Listings]
  UC3[Filter]
  UC4[Book Viewing]
  UC5[Upload Listing]
  UC6[Split Bill]
  UC7[Send Connect]
  UC8[Accept Connect]
  UC9[Chat]
  UC10[Review]
  UC13[Logout]
  UC12[Manage Users]
  Tenant --> UC1
  Tenant --> UC2
  Tenant --> UC3
  Tenant --> UC4
  Tenant --> UC6
  Tenant --> UC7
  Tenant --> UC8
  Tenant --> UC9
  Tenant --> UC10
  Tenant --> UC13
  Landlord --> UC1
  Landlord --> UC5
  Landlord --> UC2
  Landlord --> UC9
  Landlord --> UC13
  Admin --> UC12
  Admin --> UC1
  Admin --> UC13
  UC7 -.-> UC8
  UC8 -.-> UC9""",
  Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.4))

# 13 CLASS
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Class Diagram (core entities)")
add_diagram(s, """classDiagram
  class User{ +int id +String email +String role +login() +logout() +isAuthenticated() }
  class Listing{ +int id +float rent +String status +search() }
  class Appointment{ +int id +DateTime start +book() +checkConflict() }
  class ConnectionRequest{ +int id +String status +accept() +mutualAccept() }
  class Message{ +int id +send() +fetchThread() }
  class UserReview{ +int id +getAverage() }
  class BillLog{ +int id +calculateShares() }
  User "1" --> "*" Listing
  User "1" --> "*" Appointment
  User "1" --> "*" ConnectionRequest
  User "1" --> "*" Message
  User "1" --> "*" UserReview
  User "1" --> "*" BillLog
  BillLog "1" --> "*" BillLogRoommate
  Listing "1" --> "*" Appointment""",
  Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.4))

# 14 SEQUENCE BOOKING
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Sequence — Booking a Viewing")
add_diagram(s, """sequenceDiagram
  actor T as Tenant
  participant B as Browser
  participant S as Server
  participant DB as MySQL
  T->>B: Select listing & date
  B->>S: GET api/available_slots
  S->>DB: SELECT booked slots
  DB-->>S: booked[]
  S-->>B: JSON {booked}
  T->>B: Click slot
  B->>S: POST api/book_viewing
  S->>DB: BEGIN + SELECT FOR UPDATE
  S->>DB: Check conflict
  S->>DB: INSERT appointment
  S->>DB: COMMIT
  S-->>B: {success, appointment}
  B-->>T: Confirmed""",
  Inches(0.7), Inches(1.8), Inches(11.9), Inches(5.4))

# 15 SEQUENCE CONNECT CHAT
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Sequence — Connect & Chat")
add_diagram(s, """sequenceDiagram
  actor A as User A
  actor B as User B
  participant S as Server
  participant DB as MySQL
  A->>S: POST connect_request A->B
  S->>DB: INSERT PENDING
  S-->>A: Request sent
  B->>S: POST connect_request B->A
  S->>DB: UPDATE both ACCEPTED
  S-->>B: Connected!
  A->>S: POST send_message
  S->>DB: INSERT message
  A->>S: GET fetch_messages
  S->>DB: SELECT thread
  S-->>A: {messages}""",
  Inches(0.7), Inches(1.8), Inches(11.9), Inches(5.4))

# 16 SEQUENCE AUTH
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Sequence — Register & Login")
add_diagram(s, """sequenceDiagram
  actor U as User
  participant B as Browser
  participant S as Server
  participant DB as MySQL
  U->>B: Register form
  B->>S: POST /auth/register.php
  S->>DB: SELECT email?
  S->>S: password_hash(bcrypt)
  S->>DB: INSERT users
  S->>S: session_regenerate_id
  S-->>B: 302 -> dashboard
  U->>B: Login + Remember me
  B->>S: POST /auth/login.php
  S->>DB: SELECT user
  S->>S: password_verify
  opt Remember
    S->>DB: UPDATE remember_token
    S->>B: Set-Cookie
  end
  S-->>B: 302 -> dashboard""",
  Inches(0.7), Inches(1.8), Inches(11.9), Inches(5.4))

# 17 ACTIVITY BOOKING
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Activity — Booking Flow")
add_diagram(s, """flowchart TD
  Start([Start]) --> A[Select listing]
  A --> B[Pick date]
  B --> C[Fetch slots]
  C --> D{Slots available?}
  D -->|No| E[No slots] --> End1([End])
  D -->|Yes| F[Show grid]
  F --> G[Click slot]
  G --> H{Conflict?}
  H -->|Yes| J[Error] --> F
  H -->|No| K{Listing available?}
  K -->|No| L[Error] --> End2([End])
  K -->|Yes| M[INSERT PENDING]
  M --> N[COMMIT]
  N --> O[Success + ID] --> End3([End])""",
  Inches(0.7), Inches(1.8), Inches(11.9), Inches(5.4))

# 18 ACTIVITY CONNECT
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Activity — Connection Flow")
add_diagram(s, """flowchart TD
  Start([Start]) --> A[A sends request]
  A --> B[INSERT PENDING]
  B --> C[Request sent]
  C --> D{B wants connect?}
  D -->|No| G[Ignored] --> End1([End])
  D -->|Yes| I[B sends back]
  I --> J{Mutual exists?}
  J -->|No| K[INSERT PENDING] --> End2([End])
  J -->|Yes| M[UPDATE ACCEPTED]
  M --> N[Connected!]
  N --> O[Chat unlocked] --> End3([End])""",
  Inches(0.7), Inches(1.8), Inches(11.9), Inches(5.4))

# 19 ACTIVITY BILL
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Activity — Bill Split Flow")
add_diagram(s, """flowchart TD
  Start([Start]) --> A[Enter bill + amount]
  A --> B[Add roommate rows]
  B --> C{Incomes > 0?}
  C -->|No| D[Error] --> B
  C -->|Yes| E[Combined income]
  E --> F[Share = income/combined x total]
  F --> G[Show table]
  G --> H{Save?}
  H -->|No| I[Results only] --> End1([End])
  H -->|Yes| J[Auth user]
  J --> K[BEGIN]
  K --> L[INSERT bill_logs]
  L --> M[INSERT bill_log_roommates]
  M --> N[COMMIT] --> End2([End])""",
  Inches(0.7), Inches(1.8), Inches(11.9), Inches(5.4))

# 20 STATE LISTING
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "State Machine — Listing Status")
add_diagram(s, """stateDiagram-v2
  [*] --> AVAILABLE : Landlord creates
  AVAILABLE --> BOOKED : Tenant books
  AVAILABLE --> HIDDEN : Landlord hides
  AVAILABLE --> [*] : Delete
  BOOKED --> AVAILABLE : Cancel
  BOOKED --> BOOKED : New booking
  HIDDEN --> AVAILABLE : Unhide
  HIDDEN --> [*] : Delete""",
  Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.8))

# 21 STATE APPOINTMENT
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "State Machine — Appointment")
add_diagram(s, """stateDiagram-v2
  [*] --> PENDING : Books slot
  PENDING --> CONFIRMED : Landlord confirms
  PENDING --> CANCELLED : Cancel / timeout
  CONFIRMED --> CANCELLED : Either cancels
  CONFIRMED --> COMPLETED : Viewing happens
  CANCELLED --> [*]
  COMPLETED --> [*]""",
  Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.8))

# 22 STATE CONNECTION
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "State Machine — Connection")
add_diagram(s, """stateDiagram-v2
  [*] --> PENDING : A sends
  PENDING --> ACCEPTED : B sends back (mutual)
  PENDING --> REJECTED : B rejects
  PENDING --> CANCELLED : A cancels
  ACCEPTED --> DISCONNECTED : Either disconnects
  ACCEPTED --> BLOCKED : Either blocks
  DISCONNECTED --> [*]
  REJECTED --> [*]
  CANCELLED --> [*]
  BLOCKED --> [*]""",
  Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.8))

# 23 STATE SESSION
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "State Machine — User Session")
add_diagram(s, """stateDiagram-v2
  [*] --> Anonymous : App loads
  Anonymous --> Authenticating : Submit login
  Authenticating --> Authenticated : Valid
  Authenticating --> Anonymous : Invalid
  Authenticated --> Authenticated : Active session
  Authenticated --> TokenRefresh : Remember-me
  TokenRefresh --> Authenticated : Refreshed
  Authenticated --> Anonymous : Logout / expire""",
  Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.8))

# 24 COMPONENT
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Component Diagram")
add_diagram(s, """graph TB
  subgraph Pres[Presentation]
    Dash[index.php]
    AuthP[auth pages]
    Mkt[listings.php]
    Bill[expenses.php]
    Book[booking.php]
    Soc[chat connect review]
  end
  subgraph App[Application]
    Lay[layout.php]
    AC[auth.php]
    Hp[helpers.php]
    Boot[bootstrap.php]
  end
  subgraph APIx[API]
    L1[listings]
    C2[calculate]
    S3[available_slots]
    B4[book_viewing]
    C5[connect]
    M6[message]
  end
  subgraph Datax[Data]
    DB[database.php]
    MY[(MySQL)]
    FU[uploads]
  end
  Pres --> Lay
  Lay --> AC
  Mkt --> L1
  Bill --> C2
  Book --> S3
  Book --> B4
  Soc --> C5
  Soc --> M6
  L1 --> DB
  C2 --> DB
  S3 --> DB
  B4 --> DB
  C5 --> DB
  M6 --> DB
  DB --> MY""",
  Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.4))

# 25 DEPLOYMENT
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "UML"); heading(s, "Deployment Diagram")
add_diagram(s, """graph TB
  Browser[Web Browser]
  Mobile[Mobile Browser]
  Apache[Apache or PHP Server 8000]
  PHP[PHP 8.1 Runtime]
  MySQL[MySQL 8.0 3307]
  Uploads[Uploads Images]
  Browser -->|HTTP| Apache
  Mobile -->|HTTP| Apache
  Apache --> PHP
  PHP --> MySQL
  PHP --> Uploads""",
  Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.8))

# 26 MODULE HIGHLIGHTS
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Feature Deep Dive"); heading(s, "Module Highlights")
high = [
    ("Marketplace","Real-time fetch() filtering by price, room type & location; debounced 250ms."),
    ("Bill Split","Share = income / combined × total; dynamic roommate rows; optional DB save."),
    ("Booking","30-min slots 9:00–17:00; SELECT FOR UPDATE lock; 409 on conflict; transactional."),
    ("Listing Upload","JPG/PNG only, 5MB max; house-rule checkboxes; landlord-only via RBAC."),
    ("Social","Double opt-in Connect; 4s polling Chat with escapeHtml XSS guard; 1–5 reviews."),
    ("Auth","Session + remember-me cookie; bcrypt password_hash; roles tenant/landlord/admin."),
]
for i,(lead,rest) in enumerate(high):
    col = i % 3; rowi = i // 3
    l = Inches(0.6 + col*4.15); t = Inches(2.0 + rowi*2.2)
    box(s, l, t, Inches(3.9), Inches(1.9))
    bullets(s, l+Inches(0.2), t+Inches(0.2), Inches(3.5), Inches(1.5), [(lead,rest)], size=14, gap=3)

# 27 API
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Integration"); heading(s, "API Endpoints")
rows = [
    ["Method","Endpoint","Purpose"],
    ["GET","listings.php?api=listings","Filtered listings"],
    ["POST","expenses.php?api=calculate","Bill split breakdown"],
    ["GET","booking.php?api=available_slots","Booked slots"],
    ["POST","booking.php?api=book_viewing","Create appointment"],
    ["POST","connect_request.php","Send connection"],
    ["POST","send_message.php","Store message"],
    ["GET","fetch_messages.php","Poll thread"],
    ["POST","submit_review.php","Save review"],
    ["GET","get_user_reviews.php","Aggregated scores"],
]
table(s, Inches(0.6), Inches(1.9), Inches(12.1), rows, col_w=[1.2,4.5,5.3], fs=14)

# 28 SPRINTS
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Process"); heading(s, "Jira Sprint Management")
text(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5),
     [[("Scrum board · 4 sprints · 18 stories · 53 points planned, 83 delivered.", 19, TEXT, False)]])
rows = [
    ["Sprint","Goal","Points"],
    ["1","Foundation — DB, auth, layout, dashboard","23"],
    ["2","Core — marketplace & bill split","19"],
    ["3","Booking & listing upload","20"],
    ["4","Social, polish, testing, docs","21"],
]
table(s, Inches(0.6), Inches(2.3), Inches(12.1), rows, col_w=[1,9,2], fs=15)
box(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.0))
text(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(0.7),
     [[("Burndown: all sprints completed on time — 100% delivery across 8 weeks.", 18, ACCENT, True)]])

# 29 TESTING
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Quality"); heading(s, "Zephyr Test Management")
rows = [
    ["Cycle","Module","Tests","Pass"],
    ["1","Authentication","8","8"],
    ["2","Marketplace","6","6"],
    ["3","Bill Split","5","5"],
    ["4","Booking","7","7"],
    ["5","Listing Upload","5","5"],
    ["6","Social","8","8"],
    ["7","Regression","10","10"],
    ["Total","","49","49 (100%)"],
]
table(s, Inches(0.6), Inches(1.9), Inches(12.1), rows, col_w=[1,3,1,1], fs=15)

# 30 TEAM
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "People"); heading(s, "Team & Work Distribution")
team = [
    ("Dadhichi Sarker Shayon","Team Lead / Full-Stack — architecture, auth, social, UI, docs."),
    ("Shawki","Backend — marketplace, bill split, booking, listing upload."),
    ("Plabon Barua","Database / Backend — schema, seed, project setup."),
]
for i,(lead,rest) in enumerate(team):
    l = Inches(0.6 + i*4.15); t = Inches(2.0)
    box(s, l, t, Inches(3.9), Inches(2.2))
    bullets(s, l+Inches(0.2), t+Inches(0.2), Inches(3.5), Inches(1.8), [(lead,rest)], size=15, gap=4)
text(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.5),
     [
      [("Commits: Dadhichi 8 (47%) · Shawki 5 (29%) · Plabon 3 (18%) · Merge 1 (6%).", 17, MUTED, False)],
      [("Branch strategy: main (protected) -> feature branches dadhichi / shawki / plabon.", 17, MUTED, False)],
     ])

# 31 DEMO
s = prs.slides.add_slide(BLANK); bg(s)
kicker(s, "Try It"); heading(s, "Demo Accounts & Getting Started")
rows = [
    ["Role","Email","Password"],
    ["Tenant","ayesha@example.com","Roommate123!"],
    ["Landlord","rakib@example.com","Roommate123!"],
    ["Tenant","nusrat@example.com","Roommate123!"],
    ["Landlord","sajid@example.com","Roommate123!"],
    ["Tenant","tania@example.com","Roommate123!"],
]
table(s, Inches(0.6), Inches(1.9), Inches(7.4), rows, col_w=[1.5,3.5,2], fs=14)
box(s, Inches(8.4), Inches(1.9), Inches(4.3), Inches(2.6))
bullets(s, Inches(8.6), Inches(2.1), Inches(4.0), Inches(2.3), [
    ("Setup", " — import schema.sql & seed.sql, then php -S localhost:8000."),
    ("Windows", " — double-click start_server.bat, open localhost:8000."),
], size=14, gap=6)

# 32 CLOSING
s = prs.slides.add_slide(BLANK); bg(s)
text(s, Inches(2), Inches(2.4), Inches(9.3), Inches(1.0),
     [[("ROOMMATESYNC", 56, WHITE, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(2), Inches(3.6), Inches(9.3), Inches(0.8),
     [[("Find a room, split costs, manage house life — all in one place.", 22, ACCENT, False)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(2), Inches(4.7), Inches(9.3), Inches(0.5),
     [[("Built with by Team RoommateSync · ISD Course Project, KUET", 16, MUTED, False)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(2), Inches(5.5), Inches(9.3), Inches(0.6),
     [[("● PHP 8.1   ● MySQL 8.0   ● 100% Tests Passing   ● 7 Modules", 15, MUTED, False)]],
     align=PP_ALIGN.CENTER)

prs.save(PPTX_PATH)
print("Saved:", PPTX_PATH, "slides:", len(prs.slides._sldIdLst))
